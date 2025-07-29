# ppt_parser.py

from pptx import Presentation
import fitz  # PyMuPDF

def extract_text_from_ppt(file_path):
    prs = Presentation(file_path)
    all_text = []
    for slide in prs.slides:
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        all_text.append(slide_text.strip())
    return "\n\n".join(all_text)

def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

